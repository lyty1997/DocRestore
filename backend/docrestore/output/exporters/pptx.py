# Copyright 2026 @lyty1997
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""pptx 导出器（D5）：``document.md`` → PowerPoint，走 python-pptx 逐页自拼。

**为何不用 ``pandoc -t pptx``**（4 轮 spike 实证，2026-06-23）：pandoc 的
pptx 写法**无法让文字与图片同处一张 slide**（块级图片必单独成页、内联图片被丢）。
故改用 python-pptx 自拼页。

**为何按块解析 / 渲染原生表格**（2026-06-23 修复）：早期版本把整行文本直接塞进
文本框，导致 ``document.md`` 里的 HTML ``<table>`` / ``<div>`` 原始标记**当字面文本
漏到 slide 上**（用户可见一长串 ``<table border=1 ...>``）。现改为把一页拆成
**有序块**（正文 / 表格 / 图片）：散文剥掉 HTML 标签只留文本、``<table>`` 复用
公共解析层（:mod:`~docrestore.output.exporters.html_table`）渲染成**原生 pptx
表格**（含合并区）、图片保序嵌入，竖向堆叠。``$..$`` 公式留 TeX 文本（lite 不渲染）。

python-pptx 是纯 Python 依赖：**惰性导入** fail-closed
（:class:`ExportToolUnavailable`）——注册表启动导入本模块，顶层不 import pptx。
详见 ``docs/zh/export-mode.md`` §9.2。
"""

from __future__ import annotations

import html
import logging
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any, TypeAlias  # python-pptx 运行期对象无类型存根，边界处标 Any

from docrestore.output.exporters.base import (
    ExportFailed,
    ExportToolUnavailable,
)
from docrestore.output.exporters.html_table import (
    GridCells,
    Merge,
    RawCell,
    build_grid,
    grid_dimensions,
    parse_one_table,
)
from docrestore.output.ppt_layout import (
    PptLayout,
    PptLayoutPage,
    load_ppt_layout,
    region_box_emu,
)

logger = logging.getLogger(__name__)

#: 水平线行（页分隔）：``---`` / ``***`` / ``___`` 三连及以上
_HR_SPLIT = re.compile(r"(?m)^[ \t]*(?:-{3,}|\*{3,}|_{3,})[ \t]*$")
#: 顶层 H1（doc 模式无 ``---`` 时的回退切页锚）
_H1 = re.compile(r"^#\s+")
#: ATX 标题（任意级）
_ATX = re.compile(r"^(#{1,6})\s+(.*?)\s*#*\s*$")
#: 页内块级 token：HTML 表格 / markdown 图片 / HTML 图片（保序扫描）
_BLOCK_RE = re.compile(
    r"(?P<table><table\b[^>]*>.*?</table>)"
    r"|!\[[^\]]*\]\((?P<img_md>[^)]+)\)"
    r"""|<img\b[^>]*\bsrc\s*=\s*["'](?P<img_html>[^"']+)["'][^>]*>""",
    re.DOTALL | re.IGNORECASE,
)
#: ``<br>`` → 换行
_BR_RE = re.compile(r"<br\s*/?>", re.IGNORECASE)
#: 块级闭合标签 → 换行（保留分行语义）
_BLOCK_END_RE = re.compile(
    r"</(?:div|p|li|tr|h[1-6]|blockquote)\s*>", re.IGNORECASE,
)
#: 其余 HTML 标签 → 删除（保留内部文本）
_TAG_RE = re.compile(r"<[^>]+>")

# ── 版式常量（EMU；1 inch = 914400 EMU；16:9 = 13.333in × 7.5in）──────────
_EMU_PER_PT = 12700
_SLIDE_W = 12192000
_SLIDE_H = 6858000
_MARGIN = 365760        # 0.4in
_TITLE_TOP = 274320     # 0.3in
_TITLE_H = 822960       # 0.9in
_CONTENT_TOP = 1188720  # 1.3in
_GAP = 137160           # 0.15in 块间距
_TITLE_PT = 28
_BODY_PT = 14
_TABLE_PT = 11
#: 正文行高（含行距）估算：仅用于把下一块堆到合适位置
_BODY_LINE_H = int(_BODY_PT * 1.5 * _EMU_PER_PT)
#: 表格单行高估算
_TABLE_ROW_H = int(_TABLE_PT * 2.2 * _EMU_PER_PT)
#: 正文每行约可容字宽（按近似全宽字，偏保守→多留空避免重叠）
_BODY_CHAR_W = int(_BODY_PT * 0.95 * _EMU_PER_PT)
#: 内容区底边
_CONTENT_BOTTOM = _SLIDE_H - _MARGIN
#: 单块最小可用高度（剩余不足时给块的兜底高度，宁可轻微出血也不丢内容）
_MIN_BLOCK_H = _EMU_PER_PT * 24


@dataclass(frozen=True)
class _TextBlock:
    """一段正文（多行，已剥 HTML 标签）。"""

    lines: list[str]


@dataclass(frozen=True)
class _TableBlock:
    """一个 HTML 表格（已解析成行）。"""

    rows: list[list[RawCell]]


@dataclass(frozen=True)
class _ImageBlock:
    """一张图片引用（相对 src）。"""

    src: str


#: 页内有序块（正文 / 表格 / 图片）
_Block: TypeAlias = _TextBlock | _TableBlock | _ImageBlock


def _split_pages(markdown: str) -> list[str]:
    """切页：优先按水平线分隔；无分隔则回退按顶层 H1；都没有则整篇一页。"""
    pages = [p.strip() for p in _HR_SPLIT.split(markdown) if p.strip()]
    if len(pages) > 1:
        return pages
    return _split_by_headings(markdown)


def _split_by_headings(markdown: str) -> list[str]:
    """回退切页：每个顶层 ``#`` 起一页；首个标题前的引言并入第一页。"""
    lines = markdown.splitlines()
    idxs = [i for i, line in enumerate(lines) if _H1.match(line)]
    if len(idxs) <= 1:
        whole = markdown.strip()
        return [whole] if whole else []
    pages: list[str] = []
    for k, start in enumerate(idxs):
        end = idxs[k + 1] if k + 1 < len(idxs) else len(lines)
        pages.append("\n".join(lines[start:end]).strip())
    if idxs[0] > 0:
        preamble = "\n".join(lines[: idxs[0]]).strip()
        if preamble:
            pages[0] = f"{preamble}\n\n{pages[0]}"
    return [p for p in pages if p]


def _strip_html(text: str) -> str:
    """剥 HTML 标签保留内部文本：``<br>``/块级闭合 → 换行，其余标签删，实体解码。"""
    text = _BR_RE.sub("\n", text)
    text = _BLOCK_END_RE.sub("\n", text)
    text = _TAG_RE.sub("", text)
    return html.unescape(text)


def _prose_lines(prose: str) -> list[str]:
    """块 token 之间的散文 → 去标签、按行拆、丢空行。"""
    return [ln.strip() for ln in _strip_html(prose).splitlines() if ln.strip()]


def _heading_or_text(line: str, title: str | None) -> tuple[str | None, str | None]:
    """一行散文 → ``(更新后的标题, 要追加的正文行或 None)``。

    首个 ATX 标题升为 slide 标题（不进正文），其余标题降为正文行。
    """
    heading = _ATX.match(line)
    if heading is None:
        return title, line
    text = heading.group(2).strip()
    if title is None:
        return text, None
    return title, text


def _collect_prose(chunk: str, title: str | None, lines: list[str]) -> str | None:
    """把一段散文的文本行追加进 ``lines``（原地），返回更新后的标题。"""
    for line in _prose_lines(chunk):
        title, text = _heading_or_text(line, title)
        if text is not None:
            lines.append(text)
    return title


def _match_to_block(match: re.Match[str]) -> _Block | None:
    """块 token 匹配 → 表格 / 图片块；空表 / 空 src → ``None``。"""
    table_html = match.group("table")
    if table_html is not None:
        rows = parse_one_table(table_html)
        return _TableBlock(rows) if rows else None
    src = match.group("img_md") or match.group("img_html")
    return _ImageBlock(src.strip()) if src else None


def _parse_page(page: str) -> tuple[str | None, list[_Block]]:
    """拆一页 → ``(标题, 有序块列表)``。

    ``<table>`` / 图片按出现顺序成独立块，块之间的散文剥 HTML 后并入文本块——
    原始 ``<div>/<table>`` 标记不再当字面文本漏到 slide 上。
    """
    title: str | None = None
    blocks: list[_Block] = []
    text_lines: list[str] = []
    pos = 0
    for match in _BLOCK_RE.finditer(page):
        title = _collect_prose(page[pos : match.start()], title, text_lines)
        block = _match_to_block(match)
        if block is not None:
            if text_lines:
                blocks.append(_TextBlock(text_lines))
                text_lines = []
            blocks.append(block)
        pos = match.end()
    title = _collect_prose(page[pos:], title, text_lines)
    if text_lines:
        blocks.append(_TextBlock(text_lines))
    return title, blocks


def _resolve_image(doc_dir: Path, src: str) -> Path | None:
    """把图片相对引用解析为 doc_dir 内的真实文件；越界 / 远程 / 缺失 → ``None``。"""
    if "://" in src or src.startswith("data:"):
        return None
    base = doc_dir.resolve()
    candidate = (base / src).resolve()
    try:
        candidate.relative_to(base)
    except ValueError:
        return None
    return candidate if candidate.is_file() else None


def _add_title(slide: Any, title: str) -> None:
    """顶部标题框（加粗大字）。"""
    from pptx.util import Emu, Pt  # noqa: PLC0415 — pptx 已加载，懒导入保持模块轻量

    box = slide.shapes.add_textbox(
        Emu(_MARGIN), Emu(_TITLE_TOP), Emu(_SLIDE_W - 2 * _MARGIN), Emu(_TITLE_H),
    )
    frame = box.text_frame
    frame.word_wrap = True
    run = frame.paragraphs[0].add_run()
    run.text = title
    run.font.size = Pt(_TITLE_PT)
    run.font.bold = True


def _estimate_text_height(lines: list[str], width: int) -> int:
    """估算文本块高度（按近似字宽折行计行数）：仅用于把下一块堆到合适位置。"""
    per_line = max(1, width // _BODY_CHAR_W)
    wrapped = sum(max(1, -(-len(line) // per_line)) for line in lines)
    return wrapped * _BODY_LINE_H


def _add_text_block(slide: Any, lines: list[str], top: int, width: int) -> int:
    """正文文本框（每行一段，``$..$`` 公式作为普通文本），返回估算高度供堆叠。"""
    from pptx.util import Emu, Pt  # noqa: PLC0415

    height = _estimate_text_height(lines, width)
    box = slide.shapes.add_textbox(Emu(_MARGIN), Emu(top), Emu(width), Emu(height))
    frame = box.text_frame
    frame.word_wrap = True
    for i, line in enumerate(lines):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size = Pt(_BODY_PT)
    return height


def _populate_table(
    table: Any, cells: GridCells, merges: list[Merge], font_pt: int,
) -> None:
    """填充原生 pptx 表格：先逐单元格写文本 + 设字号，再做合并（块流/定位共用）。"""
    from pptx.util import Pt  # noqa: PLC0415

    for (r, c), text in cells.items():  # 先填原始单元格文本，再合并
        cell = table.cell(r, c)
        cell.text = text
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(font_pt)
    for r1, c1, r2, c2 in merges:
        table.cell(r1, c1).merge(table.cell(r2, c2))


def _add_table_block(
    slide: Any, rows: list[list[RawCell]], top: int, width: int, max_h: int,
) -> int:
    """把 HTML 表渲染成原生 pptx 表格（含合并区），返回占用高度。"""
    from pptx.util import Emu  # noqa: PLC0415

    cells, merges = build_grid(rows)
    nrows, ncols = grid_dimensions(cells, merges)
    if nrows <= 0 or ncols <= 0:
        return 0
    height = min(nrows * _TABLE_ROW_H, max(_TABLE_ROW_H, max_h))
    table = slide.shapes.add_table(
        nrows, ncols, Emu(_MARGIN), Emu(top), Emu(width), Emu(height),
    ).table
    col_w = max(1, width // ncols)
    for column in table.columns:
        column.width = Emu(col_w)
    _populate_table(table, cells, merges, _TABLE_PT)
    return height


def _add_image_block(
    slide: Any, src: str, doc_dir: Path, top: int, width: int, max_h: int,
) -> int:
    """放置一张图片（等比缩放到 ``width × max_h`` 内），返回占用高度；失败返回 0。"""
    from PIL import Image  # noqa: PLC0415 — Pillow 是硬依赖，懒导入保持模块轻量
    from pptx.util import Emu  # noqa: PLC0415

    path = _resolve_image(doc_dir, src)
    if path is None or max_h <= 0:
        return 0
    try:
        with Image.open(path) as img:
            px_w, px_h = img.size
    except (OSError, ValueError):
        return 0
    if px_w <= 0 or px_h <= 0:
        return 0
    scale = min(width / px_w, max_h / px_h)
    disp_w = max(1, int(px_w * scale))
    disp_h = max(1, int(px_h * scale))
    slide.shapes.add_picture(
        str(path), Emu(_MARGIN), Emu(top), width=Emu(disp_w), height=Emu(disp_h),
    )
    return disp_h


def _render_slide(
    slide: Any, title: str | None, blocks: list[_Block], doc_dir: Path,
) -> None:
    """把一页内容竖向堆叠到一张 slide：标题在顶，正文 / 表格 / 图片按序往下。"""
    if title:
        _add_title(slide, title)
    width = _SLIDE_W - 2 * _MARGIN
    top = _CONTENT_TOP
    for block in blocks:
        remaining = max(_CONTENT_BOTTOM - top, _MIN_BLOCK_H)  # 兜底高度防丢内容
        if isinstance(block, _TableBlock):
            used = _add_table_block(slide, block.rows, top, width, remaining)
        elif isinstance(block, _ImageBlock):
            used = _add_image_block(slide, block.src, doc_dir, top, width, remaining)
        else:
            used = _add_text_block(slide, block.lines, top, width)
        if used > 0:
            top += used + _GAP


# ── 版面定位渲染（Phase-2b）：按 .ppt_layout.json 的 bbox 摆放区域 ───────────

#: 标题类标签：定位文本框加粗（视觉层级，字号仍 lite 固定）
_LABEL_TITLE = frozenset({"paragraph_title", "figure_title"})


def _add_positioned_text(
    slide: Any, content: str, label: str, box: tuple[int, int, int, int],
) -> bool:
    """在 box（EMU）位置放文本框；剥标签后无内容则不放、返回是否真放了内容。"""
    from pptx.util import Emu, Pt  # noqa: PLC0415

    lines = _prose_lines(content)  # 去标签防御 + 按行拆
    if not lines:
        return False
    left, top, width, height = box
    frame = slide.shapes.add_textbox(
        Emu(left), Emu(top), Emu(width), Emu(height),
    ).text_frame
    frame.word_wrap = True
    bold = label in _LABEL_TITLE
    for i, line in enumerate(lines):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size = Pt(_BODY_PT)
        run.font.bold = bold
    return True


def _add_positioned_table(
    slide: Any, content: str, box: tuple[int, int, int, int],
) -> bool:
    """在 box（EMU）位置渲染原生 pptx 表格（HTML 表），返回是否成功。"""
    from pptx.util import Emu  # noqa: PLC0415

    rows = parse_one_table(content)
    if not rows:
        return False
    cells, merges = build_grid(rows)
    nrows, ncols = grid_dimensions(cells, merges)
    if nrows <= 0 or ncols <= 0:
        return False
    left, top, width, height = box
    table = slide.shapes.add_table(
        nrows, ncols, Emu(left), Emu(top), Emu(width), Emu(height),
    ).table
    col_w = max(1, width // ncols)
    for column in table.columns:
        column.width = Emu(col_w)
    _populate_table(table, cells, merges, _TABLE_PT)
    return True


def _add_positioned_image(
    slide: Any, image_ref: str, doc_dir: Path, box: tuple[int, int, int, int],
) -> bool:
    """在 box（EMU）内等比缩放居中放图片，返回是否成功。"""
    from PIL import Image  # noqa: PLC0415 — Pillow 是硬依赖，懒导入保持模块轻量
    from pptx.util import Emu  # noqa: PLC0415

    path = _resolve_image(doc_dir, image_ref)
    if path is None:
        return False
    try:
        with Image.open(path) as img:
            px_w, px_h = img.size
    except (OSError, ValueError):
        return False
    if px_w <= 0 or px_h <= 0:
        return False
    left, top, width, height = box
    scale = min(width / px_w, height / px_h)
    disp_w = max(1, int(px_w * scale))
    disp_h = max(1, int(px_h * scale))
    off_x = left + (width - disp_w) // 2  # box 内居中
    off_y = top + (height - disp_h) // 2
    slide.shapes.add_picture(
        str(path), Emu(off_x), Emu(off_y), width=Emu(disp_w), height=Emu(disp_h),
    )
    return True


def _render_positioned_page(
    slide: Any, page: PptLayoutPage, canvas: tuple[int, int], doc_dir: Path,
) -> int:
    """按 bbox 定位渲染一页区域，返回成功落下的 shape 数（0=该页无可用区域）。"""
    rendered = 0
    for region in page.regions:
        box = region_box_emu(canvas, page.image_size, region.bbox)
        if box is None:  # bbox 非法 / 零面积 → 跳过该区域
            continue
        if region.image_ref:
            added = _add_positioned_image(slide, region.image_ref, doc_dir, box)
        elif region.label == "table":
            added = _add_positioned_table(slide, region.content, box)
        else:
            added = _add_positioned_text(slide, region.content, region.label, box)
        if added:
            rendered += 1
    return rendered


def _build_block_flow(markdown: str, doc_dir: Path) -> Any:
    """竖排块流（现状 / fallback）：切页 → 每页一 slide 块级竖排，返回 Presentation。"""
    from pptx import Presentation  # noqa: PLC0415 — 惰性导入，缺依赖不阻塞启动
    from pptx.util import Emu  # noqa: PLC0415

    prs = Presentation()
    prs.slide_width = Emu(_SLIDE_W)  # Emu 是 Length 子类，满足类型
    prs.slide_height = Emu(_SLIDE_H)
    blank = prs.slide_layouts[6]  # 默认模板的 Blank 版式
    pages = _split_pages(markdown)
    for page in pages:
        slide = prs.slides.add_slide(blank)
        title, blocks = _parse_page(page)
        _render_slide(slide, title, blocks, doc_dir)
    if not pages:  # 空文档兜底：至少一张空 slide
        prs.slides.add_slide(blank)
    return prs


def _build_positioned(markdown: str, doc_dir: Path, layout: PptLayout) -> Any:
    """版面定位渲染：画布按 sidecar 尺寸，逐页按 bbox 摆区域；某页无可用区域 →
    该页退竖排（按 document.md 对应页），返回 Presentation。任一异常向上抛，由
    调用方退整篇竖排。"""
    from pptx import Presentation  # noqa: PLC0415
    from pptx.util import Emu  # noqa: PLC0415

    prs = Presentation()
    cw, ch = layout.slide_size_emu
    prs.slide_width = Emu(cw)
    prs.slide_height = Emu(ch)
    blank = prs.slide_layouts[6]
    block_pages = _split_pages(markdown)  # 供空页回退
    for i, page in enumerate(layout.pages):
        slide = prs.slides.add_slide(blank)
        rendered = _render_positioned_page(
            slide, page, layout.slide_size_emu, doc_dir,
        )
        if rendered == 0 and i < len(block_pages):  # 该页无可用区域 → 退竖排
            title, blocks = _parse_page(block_pages[i])
            _render_slide(slide, title, blocks, doc_dir)
    if not layout.pages:
        prs.slides.add_slide(blank)
    return prs


def _build_presentation(
    markdown: str, doc_dir: Path, layout: PptLayout | None,
) -> Any:
    """选渲染路径：sidecar 合法 → 定位（失败退竖排）；缺/非法 → 竖排块流。"""
    if layout is not None:
        try:
            return _build_positioned(markdown, doc_dir, layout)
        except Exception:  # noqa: BLE001 — 定位 best-effort，任一异常退整篇竖排
            logger.warning("PPT 版面定位渲染失败，退回竖排块流", exc_info=True)
    return _build_block_flow(markdown, doc_dir)


class PptxExporter:
    """``document.md`` → pptx（python-pptx，逐页一 slide，块级竖排）。"""

    suffix = "pptx"
    tool = "python-pptx"

    def ensure_available(self) -> None:
        """python-pptx 不可导入 → fail-closed。"""
        try:
            import pptx  # noqa: F401, PLC0415 — 惰性导入仅做可用性探测
        except ImportError as exc:
            raise ExportToolUnavailable(self.tool) from exc

    def export(
        self,
        doc_md: Path,
        assets_dir: Path,  # noqa: ARG002 — 图片走 document.md 相对引用解析
        out_path: Path,
    ) -> None:
        """sidecar 在 → 按 bbox 定位渲染（失败退竖排）；缺 → 竖排块流 → 保存 pptx。"""
        try:
            import pptx  # noqa: F401, PLC0415 — 惰性导入仅做可用性探测
        except ImportError as exc:
            raise ExportToolUnavailable(self.tool) from exc

        markdown = doc_md.read_text(encoding="utf-8")
        doc_dir = doc_md.parent
        out_path.parent.mkdir(parents=True, exist_ok=True)
        layout = load_ppt_layout(doc_dir)  # 缺/损坏 → None → 退竖排块流
        try:
            prs = _build_presentation(markdown, doc_dir, layout)
            prs.save(str(out_path))
        except Exception as exc:  # noqa: BLE001 — python-pptx 异常统一 fail-closed
            raise ExportFailed(self.tool, self.suffix, str(exc)[:500]) from exc
