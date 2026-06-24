# Copyright 2026 @lyty1997
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0

"""Epic D · D5 pptx 导出（python-pptx 逐页自拼）测试。

- 切页 / 页内容抽取 / 图片解析：纯函数，不依赖 python-pptx。
- 缺 python-pptx fail-closed：模拟 ImportError → ``ExportToolUnavailable``。
- 真导出（需 python-pptx）：2 页 md → pptx，python-pptx 读回 slide 数=页数、
  **从输入派生**的标题文字、图片嵌入（``<img>`` 与 ``![]()`` 都解析）、公式 TeX 文本在。
"""

from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

import pytest
from PIL import Image

from docrestore.output.exporters.base import ExportToolUnavailable
from docrestore.output.exporters.pptx import (
    PptxExporter,
    _ImageBlock,
    _parse_page,
    _resolve_image,
    _split_pages,
    _TableBlock,
    _TextBlock,
)
from docrestore.output.ppt_layout import (
    PptLayout,
    PptLayoutPage,
    PptLayoutRegion,
    region_box_emu,
    write_ppt_layout,
)


def _pptx_ready() -> bool:
    return importlib.util.find_spec("pptx") is not None


pptx_required = pytest.mark.skipif(
    not _pptx_ready(), reason="python-pptx 不可用",
)

#: 从输入派生的关键内容（构造输入）
_TITLE_A = "幻灯片标题甲ALPHA"
_TITLE_B = "幻灯片标题乙BETA"
_INLINE_TEX = "$E=mc^2$"
_CELL_TEXT = "单元格内容Zeta"
_DIV_TEXT = "居中说明文字Delta"
_MARKDOWN = f"""# {_TITLE_A}

第一页正文，行内公式 {_INLINE_TEX}。

<div style="text-align:center;">{_DIV_TEXT}</div>

<table border="1"><tr><td>列甲</td><td>列乙</td></tr>\
<tr><td>{_CELL_TEXT}</td><td>123</td></tr></table>

<img src="images/pic_0.png" alt="" />

---

# {_TITLE_B}

第二页正文。

![配图](images/pic_0.png)

$$ \\frac{{a}}{{b}} = \\sqrt{{c}} $$
"""


def _build_doc(doc_dir: Path) -> Path:
    (doc_dir / "document.md").write_text(_MARKDOWN, encoding="utf-8")
    images = doc_dir / "images"
    images.mkdir()
    Image.new("RGB", (160, 100), (200, 120, 60)).save(images / "pic_0.png")
    return doc_dir / "document.md"


class TestSplitAndExtract:
    """纯函数：切页 / 页内容抽取 / 图片解析（无依赖）。"""

    def test_split_by_horizontal_rule(self) -> None:
        pages = _split_pages("# A\n\nx\n\n---\n\n# B\n\ny")
        assert len(pages) == 2

    def test_split_fallback_by_heading(self) -> None:
        pages = _split_pages("# A标题\n\n正文a\n\n# B标题\n\n正文b")
        assert len(pages) == 2

    def test_split_single_page(self) -> None:
        assert _split_pages("只有一段正文，无分隔") == ["只有一段正文，无分隔"]

    def test_parse_page_blocks_in_order(self) -> None:
        title, blocks = _parse_page(
            "# T标题\n\n正文行ABC\n\n<div>居中XYZ</div>\n\n"
            '<table border="1"><tr><td>aa</td><td>bb</td></tr></table>\n\n'
            '<img src="images/a.png">\n\n![](images/b.png)',
        )
        assert title == "T标题"
        kinds = [type(b).__name__ for b in blocks]
        assert "_TableBlock" in kinds  # HTML 表成独立表格块
        assert kinds.count("_ImageBlock") == 2  # <img> 与 ![]() 各一块
        # div 内文保留为正文行，HTML 标签全部剥除（不再当字面文本漏出）
        text_lines = [
            ln for b in blocks if isinstance(b, _TextBlock) for ln in b.lines
        ]
        assert "正文行ABC" in text_lines
        assert "居中XYZ" in text_lines
        assert all("<" not in ln for ln in text_lines)

    def test_parse_page_table_and_image_blocks(self) -> None:
        _title, blocks = _parse_page(
            '<table border="1"><tr><td>x</td></tr></table>\n\n![](images/p.png)',
        )
        assert isinstance(blocks[0], _TableBlock)
        assert isinstance(blocks[1], _ImageBlock)
        assert blocks[1].src == "images/p.png"

    def test_parse_page_no_title(self) -> None:
        title, blocks = _parse_page("纯正文一\n\n纯正文二")
        assert title is None
        text_lines = [
            ln for b in blocks if isinstance(b, _TextBlock) for ln in b.lines
        ]
        assert text_lines == ["纯正文一", "纯正文二"]


class TestResolveImage:
    """图片引用解析：越界 / 远程 / 缺失 → None；存在 → 真实路径。"""

    def test_traversal_blocked(self, tmp_path: Path) -> None:
        assert _resolve_image(tmp_path, "../secret.png") is None

    def test_remote_skipped(self, tmp_path: Path) -> None:
        assert _resolve_image(tmp_path, "http://x/y.png") is None

    def test_missing_skipped(self, tmp_path: Path) -> None:
        assert _resolve_image(tmp_path, "images/none.png") is None

    def test_existing_resolved(self, tmp_path: Path) -> None:
        (tmp_path / "images").mkdir()
        target = tmp_path / "images" / "pic.png"
        Image.new("RGB", (8, 8), (1, 2, 3)).save(target)
        assert _resolve_image(tmp_path, "images/pic.png") == target


class TestPptxFailClosed:
    """缺 python-pptx → ensure_available fail-closed。"""

    def test_missing_pptx_raises(
        self, monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        monkeypatch.setitem(sys.modules, "pptx", None)
        with pytest.raises(ExportToolUnavailable):
            PptxExporter().ensure_available()


@pptx_required
class TestPptxExport:
    """真导出（需 python-pptx）：slide 数=页数 + 派生标题 + 图片 + 公式文本。"""

    def test_roundtrip_two_slides(self, tmp_path: Path) -> None:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        doc_md = _build_doc(tmp_path)
        out = tmp_path / ".exports" / "out.pptx"

        PptxExporter().export(doc_md, tmp_path / "images", out)

        prs = Presentation(str(out))
        slides = list(prs.slides)
        assert len(slides) == 2  # 每页一 slide

        all_text = ""
        pictures = 0
        table_cells: list[str] = []
        for slide in slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    pictures += 1
                elif shape.has_table:
                    table_cells.extend(
                        cell.text for row in shape.table.rows for cell in row.cells
                    )
                elif shape.has_text_frame:
                    all_text += shape.text_frame.text
        assert _TITLE_A in all_text  # 派生：第一页标题
        assert _TITLE_B in all_text  # 派生：第二页标题
        assert _INLINE_TEX in all_text  # 公式以 TeX 文本保留
        assert "\\frac" in all_text  # 独立公式 TeX 文本保留
        assert pictures >= 2  # <img> 与 ![]() 两张图都嵌入
        # HTML 表渲染成原生表格：派生单元格在表格里
        assert any(_CELL_TEXT in cell for cell in table_cells)
        # div 内文保留为正文，但原始 HTML 标记不再当字面文本漏到 slide 上
        assert _DIV_TEXT in all_text
        assert "<table" not in all_text
        assert "<div" not in all_text


# ── 版面定位导出（Phase-2b）：sidecar 在 → 按 bbox 定位 ──────────────────────

_POS_TITLE = "定位标题Gamma"
_POS_CELL = "定位单元格Theta"
_POS_IMG = "images/pos_0.png"
_CANVAS = (12192000, 6858000)  # 1920x1080 首页 → 16:9 画布
_IMG_SIZE = (1920, 1080)


def _positioned_doc(doc_dir: Path) -> Path:
    """构造 document.md + 图片 + .ppt_layout.json（单页：标题/表格/图片三区域）。"""
    (doc_dir / "document.md").write_text(
        f"# {_POS_TITLE}\n\n占位正文\n", encoding="utf-8",
    )
    images = doc_dir / "images"
    images.mkdir()
    Image.new("RGB", (120, 80), (40, 160, 200)).save(images / "pos_0.png")
    layout = PptLayout(
        slide_size_emu=_CANVAS,
        pages=[PptLayoutPage(
            filename="slideA.jpg",
            image_size=_IMG_SIZE,
            regions=[
                PptLayoutRegion((0, 0, 1920, 200), "paragraph_title", _POS_TITLE),
                PptLayoutRegion(
                    (0, 300, 960, 800), "table",
                    f'<table border="1"><tr><td>{_POS_CELL}</td></tr></table>',
                ),
                PptLayoutRegion(
                    (1000, 300, 1920, 800), "image", "", image_ref=_POS_IMG,
                ),
            ],
        )],
    )
    write_ppt_layout(doc_dir, layout)
    return doc_dir / "document.md"


@pptx_required
class TestPptxPositioned:
    """sidecar 在 → 按 bbox 定位渲染；非法/空页 → fail-safe 退竖排。"""

    def test_positioned_places_regions_at_bbox(self, tmp_path: Path) -> None:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        doc_md = _positioned_doc(tmp_path)
        out = tmp_path / ".exports" / "pos.pptx"
        PptxExporter().export(doc_md, tmp_path / "images", out)

        prs = Presentation(str(out))
        # 画布按 sidecar 尺寸（首页长宽比）
        assert prs.slide_width == _CANVAS[0]
        assert prs.slide_height == _CANVAS[1]
        slides = list(prs.slides)
        assert len(slides) == 1  # 一个 layout 页 → 一张 slide

        title_box = region_box_emu(_CANVAS, _IMG_SIZE, (0, 0, 1920, 200))
        table_box = region_box_emu(_CANVAS, _IMG_SIZE, (0, 300, 960, 800))
        image_box = region_box_emu(_CANVAS, _IMG_SIZE, (1000, 300, 1920, 800))
        assert title_box is not None
        assert table_box is not None
        assert image_box is not None

        title_shape = None
        table_shape = None
        picture = None
        for shape in slides[0].shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture = shape
            elif shape.has_table:
                table_shape = shape
            elif shape.has_text_frame and _POS_TITLE in shape.text_frame.text:
                title_shape = shape

        # 标题文本框落在标题区域 bbox（exporter 调 region_box_emu 定位）
        assert title_shape is not None
        assert title_shape.left == title_box[0]
        assert title_shape.top == title_box[1]
        assert title_shape.width == title_box[2]
        # 表格渲染成原生表格，落在表格区域、含派生单元格
        assert table_shape is not None
        assert table_shape.left == table_box[0]
        assert table_shape.top == table_box[1]
        cells = [c.text for row in table_shape.table.rows for c in row.cells]
        assert any(_POS_CELL in c for c in cells)
        # 图片嵌入并落在图片区域 box 内（box 内居中）
        assert picture is not None
        assert image_box[0] <= picture.left
        assert picture.left + picture.width <= image_box[0] + image_box[2]

    def test_corrupt_sidecar_falls_back_to_block_flow(
        self, tmp_path: Path,
    ) -> None:
        from pptx import Presentation

        # 两页 document.md（块流按 --- 切两页）+ 损坏 sidecar
        (tmp_path / "document.md").write_text(
            "# 页一AAA\n\n正文一\n\n---\n\n# 页二BBB\n\n正文二\n",
            encoding="utf-8",
        )
        (tmp_path / ".ppt_layout.json").write_text("{ broken", encoding="utf-8")
        out = tmp_path / ".exports" / "fb.pptx"
        PptxExporter().export(
            tmp_path / "document.md", tmp_path / "images", out,
        )

        prs = Presentation(str(out))
        # 损坏 sidecar → load 返回 None → 退竖排块流（两页两 slide），不报错
        assert len(list(prs.slides)) == 2

    def test_page_with_invalid_bbox_falls_back_per_page(
        self, tmp_path: Path,
    ) -> None:
        from pptx import Presentation

        fallback_title = "退化页CCC"
        (tmp_path / "document.md").write_text(
            f"# {fallback_title}\n\n竖排正文\n", encoding="utf-8",
        )
        # 单页 sidecar，区域 bbox 零面积（非法）→ 该页无可用区域 → 退竖排
        layout = PptLayout(
            slide_size_emu=_CANVAS,
            pages=[PptLayoutPage(
                filename="slideA.jpg",
                image_size=_IMG_SIZE,
                regions=[PptLayoutRegion((10, 10, 10, 10), "text", "不该出现")],
            )],
        )
        write_ppt_layout(tmp_path, layout)
        out = tmp_path / ".exports" / "pp.pptx"
        PptxExporter().export(
            tmp_path / "document.md", tmp_path / "images", out,
        )

        prs = Presentation(str(out))
        slides = list(prs.slides)
        assert len(slides) == 1
        all_text = "".join(
            s.text_frame.text for s in slides[0].shapes if s.has_text_frame
        )
        # 该页退竖排：document.md 对应页（按 idx 对齐）的标题出现
        assert fallback_title in all_text
